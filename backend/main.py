"""
LearnVault — FastAPI Backend (Full Admin Dashboard)
"""

from __future__ import annotations

import hashlib
import io
import os
import secrets
import shutil
from datetime import datetime, timezone
from typing import Annotated

from fastapi import (
    Depends,
    FastAPI,
    File,
    HTTPException,
    Query,
    UploadFile,
    status,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, EmailStr, Field

# ── Database imports ──────
from database import SessionLocal, engine, Base
from models import (
    User,
    Course,
    Module,
    Test,
    TestQuestion,
    StudentProgress,
    Submission,
    SupportToken,
    StudentMessage,
    Enrollment,
    Certificate,
    generate_hash,
    SECRET_KEY,
)
from cert_gen import generate_certificate
from sqlalchemy.orm import Session


# ═══════════════════════════════════════════════
#  App Initialisation
# ═══════════════════════════════════════════════

app = FastAPI(
    title="LearnVault API",
    description="Secure E-Learning Platform — Backend API",
    version="2.0.0",
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Create tables on startup
Base.metadata.create_all(bind=engine)

# ── Uploads directory ──────
import pathlib
UPLOAD_DIR = pathlib.Path(__file__).resolve().parent / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)

# Seed the single admin account
ADMIN_EMAIL = "pro@learningvault.com"
ADMIN_PASSWORD = "Test@vault26"


def _hash_password(password: str) -> str:
    return hashlib.sha256(password.encode("utf-8")).hexdigest()


def _create_token(user_id: int) -> str:
    return secrets.token_urlsafe(32)


_active_tokens: dict[str, int] = {}


def _seed_admin():
    db = SessionLocal()
    try:
        existing = db.query(User).filter(User.email == ADMIN_EMAIL).first()
        if not existing:
            admin = User(
                first_name="Admin",
                last_name="LearnVault",
                email=ADMIN_EMAIL,
                password_hash=_hash_password(ADMIN_PASSWORD),
                role="admin",
            )
            db.add(admin)
            db.commit()
    finally:
        db.close()


# Allowed email domains for student registration
ALLOWED_DOMAINS = {
    "gmail.com", "outlook.com", "hotmail.com", "yahoo.com",
    "yahoo.co.in", "live.com", "icloud.com", "protonmail.com",
    "zoho.com", "aol.com", "mail.com", "rediffmail.com",
}

# Common / easily brute-forced passwords to block
BLOCKED_PASSWORDS = {
    "abc@1234", "abcd@1234", "abc@12345", "password1!", "password@1",
    "password@123", "p@ssword1", "p@ssw0rd1", "qwerty@123",
    "welcome@1", "welcome@123", "admin@123", "admin@1234",
    "test@1234", "test@12345", "letmein@1", "changeme@1",
    "12345678a!", "1234abcd!", "abcdefg@1", "iloveyou1!",
    "sunshine@1", "princess@1", "football@1", "monkey@123",
    "dragon@123", "master@123", "trustno1!", "aa@12345678",
}


def validate_password_strength(password: str, first_name: str, last_name: str, email: str) -> str | None:
    if len(password) < 8:
        return "Password must be at least 8 characters long."
    if len(password) > 16:
        return "Password must be at most 16 characters long."
    if not any(c.isupper() for c in password):
        return "Password must contain at least one uppercase letter (A-Z)."
    if not any(c.islower() for c in password):
        return "Password must contain at least one lowercase letter (a-z)."
    if not any(c.isdigit() for c in password):
        return "Password must contain at least one number (0-9)."
    if not any(c in '!@#$%^&*()_+-=[]{}|;:,.<>?/~`' for c in password):
        return "Password must contain at least one special character (!@#$%^&* etc.)."
    pw_lower = password.lower()
    if first_name and len(first_name) >= 3 and first_name.lower() in pw_lower:
        return "Password must not contain your first name."
    if last_name and len(last_name) >= 3 and last_name.lower() in pw_lower:
        return "Password must not contain your last name."
    email_local = email.split("@")[0].lower()
    if len(email_local) >= 3 and email_local in pw_lower:
        return "Password must not contain your email username."
    if pw_lower in BLOCKED_PASSWORDS:
        return "This password is too common and easy to guess. Please choose a stronger one."
    if any(pw_lower.count(c) > len(password) // 2 for c in set(pw_lower)):
        return "Password has too many repeated characters."
    return None


# ═══════════════════════════════════════════════
#  Pydantic Schemas
# ═══════════════════════════════════════════════

class UserRole:
    student = "student"
    admin = "admin"

# ── Auth ──
class RegisterRequest(BaseModel):
    first_name: str = Field(..., min_length=1, max_length=100)
    last_name: str = Field(..., min_length=1, max_length=100)
    email: EmailStr
    password: str = Field(..., min_length=8, max_length=16)

class LoginRequest(BaseModel):
    email: EmailStr
    password: str
    required_role: str | None = None  # 'admin' or 'student' — enforced if set

class TokenResponse(BaseModel):
    access_token: str
    token_type: str = "bearer"
    role: str
    name: str

# ── Course ──
class CourseCreate(BaseModel):
    name: str = Field(..., min_length=1, max_length=200)
    category: str = Field(..., min_length=1, max_length=50)
    description: str = Field(..., min_length=1)
    author: str = Field("LearnVault", max_length=200)

class CourseUpdate(BaseModel):
    name: str | None = None
    category: str | None = None
    description: str | None = None
    author: str | None = None

class CourseOut(BaseModel):
    id: int
    name: str
    category: str
    description: str
    author: str
    created_at: datetime

# ── Module ──
class ModuleCreate(BaseModel):
    course_id: int
    title: str = Field(..., min_length=1, max_length=200)
    content_type: str = Field("video", pattern="^(video|text)$")
    url: str | None = None
    order: int = 0

class ModuleOut(BaseModel):
    id: int
    course_id: int
    title: str
    content_type: str
    url: str | None
    file_path: str | None
    order: int

# ── Test ──
class TestCreate(BaseModel):
    course_id: int
    title: str = Field(..., min_length=1, max_length=200)
    time_limit_minutes: int = Field(30, ge=1, le=300)

class TestQuestionCreate(BaseModel):
    question_type: str = Field("mcq", pattern="^(mcq|written)$")
    question: str
    option_a: str | None = None
    option_b: str | None = None
    option_c: str | None = None
    option_d: str | None = None
    correct_option: int | None = None
    marks: int = 1

class TestOut(BaseModel):
    id: int
    course_id: int
    title: str
    time_limit_minutes: int
    question_count: int
    created_at: datetime

# ── Student Management ──
class StudentUpdate(BaseModel):
    first_name: str | None = None
    last_name: str | None = None
    email: EmailStr | None = None
    start_date: datetime | None = None
    end_date: datetime | None = None

class StudentOut(BaseModel):
    id: int
    first_name: str
    last_name: str
    email: str
    start_date: datetime | None
    end_date: datetime | None
    created_at: datetime

# ── Support Token ──
class TokenCreate(BaseModel):
    subject: str = Field(..., min_length=1, max_length=300)
    body: str = Field(..., min_length=1)

class TokenReply(BaseModel):
    admin_reply: str = Field(..., min_length=1)
    status: str = Field("closed", pattern="^(open|in_progress|closed)$")

# ── Messages ──
class MessageCreate(BaseModel):
    subject: str = Field(..., min_length=1, max_length=300)
    body: str = Field(..., min_length=1)

# ── Profile ──
class ProfileUpdate(BaseModel):
    first_name: str | None = None
    last_name: str | None = None
    current_password: str | None = None
    new_password: str | None = None

# ── MCQ Submit (for tests) ──
class AnswerItem(BaseModel):
    question_id: int
    selected_option: int | None = None
    written_answer: str | None = None

class MCQSubmitRequest(BaseModel):
    course_id: int
    test_id: int
    answers: list[AnswerItem]

class MCQSubmitResponse(BaseModel):
    score: int
    total: int
    percentage: float
    record_hash: str
    submitted_at: datetime


# ═══════════════════════════════════════════════
#  DB Dependency
# ═══════════════════════════════════════════════

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()

DB = Annotated[Session, Depends(get_db)]


# ═══════════════════════════════════════════════
#  Auth Dependencies
# ═══════════════════════════════════════════════

def get_current_user(token: str = Query(...), db: Session = Depends(get_db)) -> User:
    user_id = _active_tokens.get(token)
    if user_id is None:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="Invalid or expired token")
    user = db.query(User).filter(User.id == user_id).first()
    if user is None:
        raise HTTPException(status_code=status.HTTP_401_UNAUTHORIZED, detail="User not found")
    return user

CurrentUser = Annotated[User, Depends(get_current_user)]

def require_admin(user: CurrentUser) -> User:
    if user.role != UserRole.admin:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail="Admin privileges required")
    return user

AdminUser = Annotated[User, Depends(require_admin)]


# ═══════════════════════════════════════════════
#  AUTH ROUTES
# ═══════════════════════════════════════════════

@app.post("/auth/register", response_model=TokenResponse, status_code=201)
def register(body: RegisterRequest, db: DB):
    domain = body.email.split("@")[-1].lower()
    if domain not in ALLOWED_DOMAINS:
        raise HTTPException(status_code=400, detail=f"Registration is only allowed with common email providers (Gmail, Outlook, Yahoo, etc.). '{domain}' is not accepted.")
    pw_error = validate_password_strength(body.password, body.first_name, body.last_name, body.email)
    if pw_error:
        raise HTTPException(status_code=400, detail=pw_error)
    existing = db.query(User).filter(User.email == body.email).first()
    if existing:
        raise HTTPException(status_code=409, detail="An account with this email already exists. Please sign in instead.")
    user = User(
        first_name=body.first_name, last_name=body.last_name,
        email=body.email, password_hash=_hash_password(body.password), role=UserRole.student,
    )
    db.add(user)
    db.commit()
    db.refresh(user)
    token = _create_token(user.id)
    _active_tokens[token] = user.id
    return TokenResponse(access_token=token, role=user.role, name=f"{user.first_name} {user.last_name}")


@app.post("/auth/login", response_model=TokenResponse)
def login(body: LoginRequest, db: DB):
    user = db.query(User).filter(User.email == body.email).first()
    if not user:
        raise HTTPException(status_code=401, detail="No account found with this email. Please check your email or sign up first.")
    if user.password_hash != _hash_password(body.password):
        raise HTTPException(status_code=401, detail="Incorrect password. Please try again.")
    # Role-based login separation
    if body.required_role:
        if body.required_role == "student" and user.role == "admin":
            raise HTTPException(status_code=403, detail="Admin accounts cannot sign in on the student portal. Please use the admin login page.")
        if body.required_role == "admin" and user.role != "admin":
            raise HTTPException(status_code=403, detail="This login page is for administrators only. Please use the student portal.")
    token = _create_token(user.id)
    _active_tokens[token] = user.id
    return TokenResponse(access_token=token, role=user.role, name=f"{user.first_name} {user.last_name}")


@app.post("/auth/logout")
def logout(token: str = Query(...)):
    _active_tokens.pop(token, None)
    return {"detail": "Logged out"}


# ═══════════════════════════════════════════════
#  COURSE CRUD (Admin)
# ═══════════════════════════════════════════════

@app.get("/admin/courses")
def list_courses(admin: AdminUser, db: DB):
    courses = db.query(Course).order_by(Course.created_at.desc()).all()
    return [CourseOut(id=c.id, name=c.name, category=c.category, description=c.description, author=c.author, created_at=c.created_at) for c in courses]


@app.post("/admin/course/", response_model=CourseOut, status_code=201)
def create_course(body: CourseCreate, admin: AdminUser, db: DB):
    course = Course(name=body.name, category=body.category, description=body.description, author=body.author)
    db.add(course)
    db.commit()
    db.refresh(course)
    return CourseOut(id=course.id, name=course.name, category=course.category, description=course.description, author=course.author, created_at=course.created_at)


@app.put("/admin/course/{course_id}", response_model=CourseOut)
def update_course(course_id: int, body: CourseUpdate, admin: AdminUser, db: DB):
    course = db.query(Course).filter(Course.id == course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    if body.name is not None: course.name = body.name
    if body.category is not None: course.category = body.category
    if body.description is not None: course.description = body.description
    if body.author is not None: course.author = body.author
    db.commit()
    db.refresh(course)
    return CourseOut(id=course.id, name=course.name, category=course.category, description=course.description, author=course.author, created_at=course.created_at)


@app.delete("/admin/course/{course_id}")
def delete_course(course_id: int, admin: AdminUser, db: DB):
    course = db.query(Course).filter(Course.id == course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    db.delete(course)
    db.commit()
    return {"detail": f"Course '{course.name}' deleted"}


# ═══════════════════════════════════════════════
#  MODULE CRUD (Admin)
# ═══════════════════════════════════════════════

@app.post("/admin/module/", response_model=ModuleOut, status_code=201)
def create_module(body: ModuleCreate, admin: AdminUser, db: DB):
    course = db.query(Course).filter(Course.id == body.course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    module = Module(course_id=body.course_id, title=body.title, content_type=body.content_type, url=body.url, order=body.order)
    db.add(module)
    db.commit()
    db.refresh(module)
    return ModuleOut(id=module.id, course_id=module.course_id, title=module.title, content_type=module.content_type, url=module.url, file_path=module.file_path, order=module.order)


@app.put("/admin/module/{module_id}", response_model=ModuleOut)
def update_module(module_id: int, body: ModuleCreate, admin: AdminUser, db: DB):
    module = db.query(Module).filter(Module.id == module_id).first()
    if not module:
        raise HTTPException(status_code=404, detail="Module not found")
    module.title = body.title
    module.content_type = body.content_type
    module.url = body.url
    module.order = body.order
    db.commit()
    db.refresh(module)
    return ModuleOut(id=module.id, course_id=module.course_id, title=module.title, content_type=module.content_type, url=module.url, file_path=module.file_path, order=module.order)


@app.delete("/admin/module/{module_id}")
def delete_module(module_id: int, admin: AdminUser, db: DB):
    module = db.query(Module).filter(Module.id == module_id).first()
    if not module:
        raise HTTPException(status_code=404, detail="Module not found")
    db.delete(module)
    db.commit()
    return {"detail": f"Module '{module.title}' deleted"}


# ═══════════════════════════════════════════════
#  FILE UPLOAD
# ═══════════════════════════════════════════════

@app.post("/admin/upload")
def upload_file(module_id: int = Query(...), admin: AdminUser = None, file: UploadFile = File(...), db: Session = Depends(get_db)):
    module = db.query(Module).filter(Module.id == module_id).first()
    if not module:
        raise HTTPException(status_code=404, detail="Module not found")
    safe_name = f"{module_id}_{file.filename.replace(' ', '_')}"
    dest = UPLOAD_DIR / safe_name
    with open(dest, "wb") as f:
        shutil.copyfileobj(file.file, f)
    module.file_path = f"/uploads/{safe_name}"
    if not module.url:
        module.url = module.file_path
    db.commit()
    return {"file_path": module.file_path, "filename": file.filename}


@app.get("/admin/course/{course_id}/modules")
def admin_course_modules(course_id: int, admin: AdminUser, db: DB):
    """Admin endpoint to view modules for a course (no enrollment needed)."""
    course = db.query(Course).filter(Course.id == course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    modules = db.query(Module).filter(Module.course_id == course_id).order_by(Module.order).all()
    return {
        "course": CourseOut(id=course.id, name=course.name, category=course.category, description=course.description, author=course.author, created_at=course.created_at),
        "modules": [ModuleOut(id=m.id, course_id=m.course_id, title=m.title, content_type=m.content_type, url=m.url, file_path=m.file_path, order=m.order) for m in modules],
    }


# ═══════════════════════════════════════════════
#  TEST CRUD (Admin)
# ═══════════════════════════════════════════════

@app.post("/admin/test/", status_code=201)
def create_test(body: TestCreate, admin: AdminUser, db: DB):
    course = db.query(Course).filter(Course.id == body.course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    test = Test(course_id=body.course_id, title=body.title, time_limit_minutes=body.time_limit_minutes)
    db.add(test)
    db.commit()
    db.refresh(test)
    return {"id": test.id, "course_id": test.course_id, "title": test.title, "time_limit_minutes": test.time_limit_minutes}


@app.post("/admin/test/{test_id}/question", status_code=201)
def add_test_question(test_id: int, body: TestQuestionCreate, admin: AdminUser, db: DB):
    test = db.query(Test).filter(Test.id == test_id).first()
    if not test:
        raise HTTPException(status_code=404, detail="Test not found")
    q = TestQuestion(
        test_id=test_id, question_type=body.question_type, question=body.question,
        option_a=body.option_a, option_b=body.option_b, option_c=body.option_c, option_d=body.option_d,
        correct_option=body.correct_option, marks=body.marks,
    )
    db.add(q)
    db.commit()
    db.refresh(q)
    return {"id": q.id, "question": q.question, "question_type": q.question_type}


@app.get("/admin/tests")
def list_tests(admin: AdminUser, db: DB):
    tests = db.query(Test).order_by(Test.created_at.desc()).all()
    result = []
    for t in tests:
        qcount = db.query(TestQuestion).filter(TestQuestion.test_id == t.id).count()
        result.append(TestOut(id=t.id, course_id=t.course_id, title=t.title, time_limit_minutes=t.time_limit_minutes, question_count=qcount, created_at=t.created_at))
    return result


@app.delete("/admin/test/{test_id}")
def delete_test(test_id: int, admin: AdminUser, db: DB):
    test = db.query(Test).filter(Test.id == test_id).first()
    if not test:
        raise HTTPException(status_code=404, detail="Test not found")
    db.delete(test)
    db.commit()
    return {"detail": f"Test '{test.title}' deleted"}


# ═══════════════════════════════════════════════
#  STUDENT ROUTES
# ═══════════════════════════════════════════════

@app.get("/student/courses")
def student_courses(user: CurrentUser, db: DB):
    courses = db.query(Course).order_by(Course.created_at.desc()).all()
    return [CourseOut(id=c.id, name=c.name, category=c.category, description=c.description, author=c.author, created_at=c.created_at) for c in courses]


@app.get("/student/modules/{course_id}")
def student_modules(course_id: int, user: CurrentUser, db: DB):
    course = db.query(Course).filter(Course.id == course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    modules = db.query(Module).filter(Module.course_id == course_id).order_by(Module.order).all()
    tests = db.query(Test).filter(Test.course_id == course_id).all()
    # Check progress
    completed_ids = set()
    rows = db.query(StudentProgress).filter(StudentProgress.user_id == user.id, StudentProgress.course_id == course_id).all()
    for r in rows:
        completed_ids.add(r.module_id)
    total_modules = len(modules)
    all_done = total_modules > 0 and len(completed_ids) >= total_modules
    return {
        "course": CourseOut(id=course.id, name=course.name, category=course.category, description=course.description, author=course.author, created_at=course.created_at),
        "modules": [ModuleOut(id=m.id, course_id=m.course_id, title=m.title, content_type=m.content_type, url=m.url, file_path=m.file_path, order=m.order) for m in modules],
        "completed_module_ids": list(completed_ids),
        "tests": [{"id": t.id, "title": t.title, "time_limit_minutes": t.time_limit_minutes, "locked": not all_done} for t in tests],
    }


@app.post("/student/progress")
def mark_progress(course_id: int = Query(...), module_id: int = Query(...), user: User = Depends(get_current_user), db: Session = Depends(get_db)):
    # Check enrollment
    enrollment = db.query(Enrollment).filter(Enrollment.user_id == user.id, Enrollment.course_id == course_id).first()
    if not enrollment:
        raise HTTPException(status_code=403, detail="Not enrolled in this course")
    # Validate module belongs to course
    module = db.query(Module).filter(Module.id == module_id, Module.course_id == course_id).first()
    if not module:
        raise HTTPException(status_code=404, detail="Module not found in this course")
    existing = db.query(StudentProgress).filter(
        StudentProgress.user_id == user.id, StudentProgress.course_id == course_id, StudentProgress.module_id == module_id
    ).first()
    if existing:
        return {"detail": "Already marked"}
    sp = StudentProgress(user_id=user.id, course_id=course_id, module_id=module_id)
    db.add(sp)
    db.commit()
    return {"detail": "Module marked as complete"}


@app.get("/student/test/{test_id}")
def get_test(test_id: int, user: CurrentUser, db: DB):
    test = db.query(Test).filter(Test.id == test_id).first()
    if not test:
        raise HTTPException(status_code=404, detail="Test not found")
    # Check enrollment
    enrollment = db.query(Enrollment).filter(Enrollment.user_id == user.id, Enrollment.course_id == test.course_id).first()
    if not enrollment:
        raise HTTPException(status_code=403, detail="Not enrolled in this course")
    # Check if all modules are complete
    modules = db.query(Module).filter(Module.course_id == test.course_id).all()
    completed = db.query(StudentProgress).filter(
        StudentProgress.user_id == user.id, StudentProgress.course_id == test.course_id
    ).count()
    if len(modules) > 0 and completed < len(modules):
        raise HTTPException(status_code=403, detail="Complete all modules before taking the test.")
    questions = db.query(TestQuestion).filter(TestQuestion.test_id == test_id).all()
    return {
        "test": {"id": test.id, "title": test.title, "time_limit_minutes": test.time_limit_minutes, "course_id": test.course_id},
        "questions": [
            {
                "id": q.id, "question_type": q.question_type, "question": q.question, "marks": q.marks,
                "options": [q.option_a, q.option_b, q.option_c, q.option_d] if q.question_type == "mcq" else None,
            } for q in questions
        ],
    }


@app.post("/submit_mcq/", response_model=MCQSubmitResponse)
def submit_test(body: MCQSubmitRequest, user: CurrentUser, db: DB):
    test = db.query(Test).filter(Test.id == body.test_id).first()
    if not test:
        raise HTTPException(status_code=404, detail="Test not found")
    # Check enrollment
    enrollment = db.query(Enrollment).filter(Enrollment.user_id == user.id, Enrollment.course_id == body.course_id).first()
    if not enrollment:
        raise HTTPException(status_code=403, detail="Not enrolled in this course")
    # Check attempt count (max 3)
    prev_attempts = db.query(Submission).filter(
        Submission.user_id == user.id, Submission.test_id == body.test_id
    ).count()
    if prev_attempts >= 3:
        raise HTTPException(status_code=400, detail="Maximum 3 attempts reached for this test.")
    questions = db.query(TestQuestion).filter(TestQuestion.test_id == body.test_id).all()
    q_map = {q.id: q for q in questions}
    score = 0
    total = 0
    for ans in body.answers:
        q = q_map.get(ans.question_id)
        if not q:
            continue
        total += q.marks
        if q.question_type == "mcq" and ans.selected_option == q.correct_option:
            score += q.marks
    percentage = round((score / total * 100) if total > 0 else 0.0, 2)
    sub = Submission(
        user_id=user.id, course_id=body.course_id, test_id=body.test_id,
        score=score, total=total, percentage=percentage,
        attempt_number=prev_attempts + 1,
    )
    db.add(sub)
    db.commit()
    db.refresh(sub)
    return MCQSubmitResponse(score=score, total=total, percentage=percentage, record_hash=sub.record_hash, submitted_at=sub.submitted_at)


# ═══════════════════════════════════════════════
#  STUDENT MANAGEMENT (Admin)
# ═══════════════════════════════════════════════

@app.get("/admin/students")
def list_students(admin: AdminUser, db: DB):
    students = db.query(User).filter(User.role == "student").order_by(User.created_at.desc()).all()
    return [StudentOut(id=s.id, first_name=s.first_name, last_name=s.last_name, email=s.email, start_date=s.start_date, end_date=s.end_date, created_at=s.created_at) for s in students]


@app.put("/admin/student/{student_id}", response_model=StudentOut)
def update_student(student_id: int, body: StudentUpdate, admin: AdminUser, db: DB):
    student = db.query(User).filter(User.id == student_id, User.role == "student").first()
    if not student:
        raise HTTPException(status_code=404, detail="Student not found")
    if body.first_name is not None: student.first_name = body.first_name
    if body.last_name is not None: student.last_name = body.last_name
    if body.email is not None: student.email = body.email
    if body.start_date is not None: student.start_date = body.start_date
    if body.end_date is not None: student.end_date = body.end_date
    db.commit()
    db.refresh(student)
    return StudentOut(id=student.id, first_name=student.first_name, last_name=student.last_name, email=student.email, start_date=student.start_date, end_date=student.end_date, created_at=student.created_at)


@app.delete("/admin/student/{student_id}")
def delete_student(student_id: int, admin: AdminUser, db: DB):
    student = db.query(User).filter(User.id == student_id, User.role == "student").first()
    if not student:
        raise HTTPException(status_code=404, detail="Student not found")
    # Clean token cache
    for tok, uid in list(_active_tokens.items()):
        if uid == student.id:
            del _active_tokens[tok]
    db.delete(student)
    db.commit()
    return {"detail": f"Student '{student.first_name} {student.last_name}' removed"}


# ═══════════════════════════════════════════════
#  STUDENT PROGRESS (Admin)
# ═══════════════════════════════════════════════

@app.get("/admin/progress")
def admin_progress(admin: AdminUser, db: DB):
    students = db.query(User).filter(User.role == "student").all()
    result = []
    for s in students:
        courses_progress = []
        courses = db.query(Course).all()
        for c in courses:
            total_modules = db.query(Module).filter(Module.course_id == c.id).count()
            completed = db.query(StudentProgress).filter(StudentProgress.user_id == s.id, StudentProgress.course_id == c.id).count()
            submissions = db.query(Submission).filter(Submission.user_id == s.id, Submission.course_id == c.id).all()
            best_score = max((sub.percentage for sub in submissions), default=None)
            if total_modules > 0 or completed > 0 or submissions:
                courses_progress.append({
                    "course_id": c.id, "course_name": c.name,
                    "total_modules": total_modules, "completed_modules": completed,
                    "completion_pct": round(completed / total_modules * 100, 1) if total_modules > 0 else 0,
                    "best_score": best_score,
                })
        if courses_progress:
            result.append({
                "student_id": s.id, "student_name": f"{s.first_name} {s.last_name}", "email": s.email,
                "courses": courses_progress,
            })
    return result


@app.post("/admin/alert/{student_id}")
def send_alert(student_id: int, body: MessageCreate, admin: AdminUser, db: DB):
    student = db.query(User).filter(User.id == student_id, User.role == "student").first()
    if not student:
        raise HTTPException(status_code=404, detail="Student not found")
    msg = StudentMessage(student_id=student_id, subject=body.subject, body=body.body)
    db.add(msg)
    db.commit()
    return {"detail": f"Message sent to {student.first_name} {student.last_name}"}


# ═══════════════════════════════════════════════
#  SUPPORT TOKENS
# ═══════════════════════════════════════════════

@app.post("/student/token", status_code=201)
def create_token_ticket(body: TokenCreate, user: CurrentUser, db: DB):
    ticket = SupportToken(user_id=user.id, subject=body.subject, body=body.body)
    db.add(ticket)
    db.commit()
    db.refresh(ticket)
    return {"id": ticket.id, "subject": ticket.subject, "status": ticket.status}


@app.get("/student/tokens")
def student_tokens(user: CurrentUser, db: DB):
    tickets = db.query(SupportToken).filter(SupportToken.user_id == user.id).order_by(SupportToken.created_at.desc()).all()
    return [{"id": t.id, "subject": t.subject, "body": t.body, "status": t.status, "admin_reply": t.admin_reply, "created_at": t.created_at} for t in tickets]


@app.get("/student/messages")
def student_messages(user: CurrentUser, db: DB):
    msgs = db.query(StudentMessage).filter(StudentMessage.student_id == user.id).order_by(StudentMessage.created_at.desc()).all()
    return [{"id": m.id, "subject": m.subject, "body": m.body, "is_read": m.is_read, "created_at": m.created_at} for m in msgs]


@app.get("/admin/tokens")
def admin_tokens(admin: AdminUser, db: DB):
    tickets = db.query(SupportToken).order_by(SupportToken.created_at.desc()).all()
    result = []
    for t in tickets:
        user = db.query(User).filter(User.id == t.user_id).first()
        result.append({
            "id": t.id, "student_name": f"{user.first_name} {user.last_name}" if user else "Unknown",
            "student_email": user.email if user else "", "subject": t.subject, "body": t.body,
            "status": t.status, "admin_reply": t.admin_reply, "created_at": t.created_at,
        })
    return result


@app.put("/admin/token/{token_id}")
def reply_token(token_id: int, body: TokenReply, admin: AdminUser, db: DB):
    ticket = db.query(SupportToken).filter(SupportToken.id == token_id).first()
    if not ticket:
        raise HTTPException(status_code=404, detail="Ticket not found")
    if ticket.status == "closed":
        raise HTTPException(status_code=400, detail="This ticket is closed and cannot be reopened.")
    ticket.admin_reply = body.admin_reply
    ticket.status = body.status
    db.commit()
    return {"detail": "Reply saved", "status": ticket.status}


# ═══════════════════════════════════════════════
#  ADMIN PROFILE
# ═══════════════════════════════════════════════

@app.get("/admin/profile")
def get_profile(admin: AdminUser, db: DB):
    return {"id": admin.id, "first_name": admin.first_name, "last_name": admin.last_name, "email": admin.email, "role": admin.role}


@app.put("/admin/profile")
def update_profile(body: ProfileUpdate, admin: AdminUser, db: DB):
    if body.first_name is not None: admin.first_name = body.first_name
    if body.last_name is not None: admin.last_name = body.last_name
    if body.new_password:
        if not body.current_password or admin.password_hash != _hash_password(body.current_password):
            raise HTTPException(status_code=400, detail="Current password is incorrect")
        admin.password_hash = _hash_password(body.new_password)
    db.commit()
    return {"detail": "Profile updated", "first_name": admin.first_name, "last_name": admin.last_name}


# ═══════════════════════════════════════════════
#  EXCEL EXPORTS (Admin)
# ═══════════════════════════════════════════════

@app.get("/admin/export/progress")
def export_progress(admin: AdminUser, db: DB):
    import pandas as pd
    students = db.query(User).filter(User.role == "student").all()
    rows = []
    for s in students:
        for c in db.query(Course).all():
            total = db.query(Module).filter(Module.course_id == c.id).count()
            done = db.query(StudentProgress).filter(StudentProgress.user_id == s.id, StudentProgress.course_id == c.id).count()
            subs = db.query(Submission).filter(Submission.user_id == s.id, Submission.course_id == c.id).all()
            best = max((sub.percentage for sub in subs), default=None)
            if total > 0 or done > 0 or subs:
                rows.append({
                    "Student": f"{s.first_name} {s.last_name}", "Email": s.email,
                    "Course": c.name, "Modules Done": done, "Total Modules": total,
                    "Completion %": round(done / total * 100, 1) if total > 0 else 0,
                    "Best Score %": best,
                })
    df = pd.DataFrame(rows) if rows else pd.DataFrame(columns=["Student", "Email", "Course", "Modules Done", "Total Modules", "Completion %", "Best Score %"])
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as w:
        df.to_excel(w, index=False, sheet_name="Progress")
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                             headers={"Content-Disposition": "attachment; filename=learnvault-progress.xlsx"})


@app.get("/admin/export/students")
def export_students(admin: AdminUser, db: DB):
    import pandas as pd
    students = db.query(User).filter(User.role == "student").all()
    rows = [{"ID": s.id, "First Name": s.first_name, "Last Name": s.last_name, "Email": s.email,
             "Start Date": str(s.start_date) if s.start_date else "", "End Date": str(s.end_date) if s.end_date else "",
             "Registered": str(s.created_at)} for s in students]
    df = pd.DataFrame(rows) if rows else pd.DataFrame(columns=["ID", "First Name", "Last Name", "Email", "Start Date", "End Date", "Registered"])
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as w:
        df.to_excel(w, index=False, sheet_name="Students")
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                             headers={"Content-Disposition": "attachment; filename=learnvault-students.xlsx"})


@app.get("/admin/export/scores")
def export_scores(admin: AdminUser, db: DB):
    import pandas as pd
    subs = db.query(Submission).all()
    rows = []
    for s in subs:
        user = db.query(User).filter(User.id == s.user_id).first()
        course = db.query(Course).filter(Course.id == s.course_id).first()
        rows.append({
            "Student": f"{user.first_name} {user.last_name}" if user else "Unknown",
            "Course": course.name if course else "Unknown",
            "Score": s.score, "Total": s.total, "Percentage": s.percentage,
            "Hash": s.record_hash, "Date": str(s.submitted_at),
        })
    df = pd.DataFrame(rows) if rows else pd.DataFrame(columns=["Student", "Course", "Score", "Total", "Percentage", "Hash", "Date"])
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as w:
        df.to_excel(w, index=False, sheet_name="Scores")
    buf.seek(0)
    return StreamingResponse(buf, media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                             headers={"Content-Disposition": "attachment; filename=learnvault-scores.xlsx"})


# ═══════════════════════════════════════════════
#  STUDENT DASHBOARD ROUTES
# ═══════════════════════════════════════════════

@app.get("/student/catalog")
def course_catalog(user: CurrentUser, db: DB):
    """All courses the student has NOT yet enrolled in."""
    enrolled_ids = [e.course_id for e in db.query(Enrollment).filter(Enrollment.user_id == user.id).all()]
    courses = db.query(Course).filter(~Course.id.in_(enrolled_ids) if enrolled_ids else True).all()
    return [
        {
            "id": c.id, "name": c.name, "category": c.category,
            "description": c.description, "author": c.author, "created_at": c.created_at,
            "module_count": db.query(Module).filter(Module.course_id == c.id).count(),
            "test_count": db.query(Test).filter(Test.course_id == c.id).count(),
        }
        for c in courses
    ]


@app.post("/student/enroll/{course_id}", status_code=201)
def enroll_in_course(course_id: int, user: CurrentUser, db: DB):
    course = db.query(Course).filter(Course.id == course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    existing = db.query(Enrollment).filter(
        Enrollment.user_id == user.id, Enrollment.course_id == course_id
    ).first()
    if existing:
        raise HTTPException(status_code=400, detail="Already enrolled in this course")
    enrollment = Enrollment(user_id=user.id, course_id=course_id)
    db.add(enrollment)
    db.commit()
    return {"detail": f"Enrolled in {course.name}", "course_id": course_id}


@app.get("/student/enrolled")
def enrolled_courses(user: CurrentUser, db: DB):
    """Courses student is enrolled in, with progress info."""
    enrollments = db.query(Enrollment).filter(Enrollment.user_id == user.id).order_by(Enrollment.enrolled_at.desc()).all()
    result = []
    for e in enrollments:
        course = db.query(Course).filter(Course.id == e.course_id).first()
        if not course:
            continue
        modules = db.query(Module).filter(Module.course_id == course.id).order_by(Module.order).all()
        completed_ids = [
            sp.module_id for sp in db.query(StudentProgress).filter(
                StudentProgress.user_id == user.id, StudentProgress.course_id == course.id
            ).all()
        ]
        tests = db.query(Test).filter(Test.course_id == course.id).all()
        # Get best submission for each test
        test_info = []
        for t in tests:
            subs = db.query(Submission).filter(
                Submission.user_id == user.id, Submission.test_id == t.id
            ).order_by(Submission.percentage.desc()).all()
            test_info.append({
                "id": t.id, "title": t.title, "time_limit_minutes": t.time_limit_minutes,
                "attempts_used": len(subs), "max_attempts": 3,
                "best_score": subs[0].percentage if subs else None,
            })
        all_modules_done = len(completed_ids) >= len(modules) if modules else False
        # Check if certificate exists
        cert = db.query(Certificate).filter(
            Certificate.user_id == user.id, Certificate.course_id == course.id
        ).first()
        result.append({
            "course_id": course.id, "name": course.name, "category": course.category,
            "description": course.description, "author": course.author,
            "enrolled_at": e.enrolled_at,
            "modules": [
                {
                    "id": m.id, "title": m.title, "content_type": m.content_type,
                    "url": m.url, "file_path": m.file_path, "order": m.order,
                    "completed": m.id in completed_ids,
                }
                for m in modules
            ],
            "total_modules": len(modules),
            "completed_modules": len(completed_ids),
            "completion_pct": round(len(completed_ids) / len(modules) * 100, 1) if modules else 0,
            "all_modules_done": all_modules_done,
            "tests": test_info,
            "has_certificate": cert is not None,
            "certificate_id": cert.id if cert else None,
        })
    return result


@app.get("/student/course/{course_id}/modules")
def get_course_modules(course_id: int, user: CurrentUser, db: DB):
    """Get modules for a specific course with completion status."""
    enrollment = db.query(Enrollment).filter(
        Enrollment.user_id == user.id, Enrollment.course_id == course_id
    ).first()
    if not enrollment:
        raise HTTPException(status_code=403, detail="Not enrolled in this course")
    modules = db.query(Module).filter(Module.course_id == course_id).order_by(Module.order).all()
    completed_ids = [
        sp.module_id for sp in db.query(StudentProgress).filter(
            StudentProgress.user_id == user.id, StudentProgress.course_id == course_id
        ).all()
    ]
    return [
        {
            "id": m.id, "title": m.title, "content_type": m.content_type,
            "url": m.url, "file_path": m.file_path, "order": m.order,
            "completed": m.id in completed_ids,
        }
        for m in modules
    ]


@app.post("/student/certificate/{course_id}", status_code=201)
def generate_cert(course_id: int, user: CurrentUser, db: DB):
    """Generate a certificate after test completion."""
    course = db.query(Course).filter(Course.id == course_id).first()
    if not course:
        raise HTTPException(status_code=404, detail="Course not found")
    # Check enrollment
    enrollment = db.query(Enrollment).filter(
        Enrollment.user_id == user.id, Enrollment.course_id == course_id
    ).first()
    if not enrollment:
        raise HTTPException(status_code=403, detail="Not enrolled")
    # Check existing certificate
    existing = db.query(Certificate).filter(
        Certificate.user_id == user.id, Certificate.course_id == course_id
    ).first()
    if existing:
        return {"detail": "Certificate already exists", "id": existing.id, "file_path": existing.file_path}
    # Check test completion
    tests = db.query(Test).filter(Test.course_id == course_id).all()
    if not tests:
        raise HTTPException(status_code=400, detail="No tests available for this course")
    best_sub = db.query(Submission).filter(
        Submission.user_id == user.id, Submission.course_id == course_id
    ).order_by(Submission.percentage.desc()).first()
    if not best_sub:
        raise HTTPException(status_code=400, detail="You have not taken any tests yet")
    if best_sub.percentage < 70.0:
        raise HTTPException(status_code=400, detail=f"Minimum 70% required to earn a certificate. Your best score is {best_sub.percentage}%.")
    # Generate PDF
    student_name = f"{user.first_name} {user.last_name}"
    start_date = enrollment.enrolled_at.strftime("%Y-%m-%d")
    completion_date = best_sub.submitted_at.strftime("%Y-%m-%d")
    # Create cert record first to get ID
    cert = Certificate(
        user_id=user.id, course_id=course_id,
        score_percentage=best_sub.percentage, file_path="pending",
    )
    db.add(cert)
    db.commit()
    db.refresh(cert)
    # Generate PDF with cert ID
    file_path = generate_certificate(
        student_name=student_name, course_name=course.name,
        score_pct=best_sub.percentage, start_date=start_date,
        completion_date=completion_date, cert_id=cert.id,
    )
    cert.file_path = file_path
    db.commit()
    return {"detail": "Certificate generated", "id": cert.id, "file_path": file_path}


@app.get("/student/certificates")
def list_certificates(user: CurrentUser, db: DB):
    certs = db.query(Certificate).filter(Certificate.user_id == user.id).order_by(Certificate.issued_at.desc()).all()
    result = []
    for c in certs:
        course = db.query(Course).filter(Course.id == c.course_id).first()
        result.append({
            "id": c.id, "course_name": course.name if course else "Unknown",
            "score_percentage": c.score_percentage, "file_path": c.file_path,
            "issued_at": c.issued_at,
        })
    return result


@app.delete("/student/account")
def delete_own_account(user: CurrentUser, db: DB):
    """Student deletes their own account and all related data."""
    if user.role == "admin":
        raise HTTPException(status_code=400, detail="Admin account cannot be deleted")
    # Remove token
    for tok, uid in list(_active_tokens.items()):
        if uid == user.id:
            del _active_tokens[tok]
    db.delete(user)
    db.commit()
    return {"detail": "Account deleted successfully"}


@app.get("/student/profile")
def student_profile(user: CurrentUser, db: DB):
    enrolled_count = db.query(Enrollment).filter(Enrollment.user_id == user.id).count()
    cert_count = db.query(Certificate).filter(Certificate.user_id == user.id).count()
    return {
        "id": user.id, "first_name": user.first_name, "last_name": user.last_name,
        "email": user.email, "created_at": user.created_at,
        "enrolled_courses": enrolled_count, "certificates_earned": cert_count,
    }


@app.put("/student/profile")
def update_student_profile(user: CurrentUser, body: dict, db: DB):
    if body.get("first_name"):
        user.first_name = body["first_name"]
    if body.get("last_name"):
        user.last_name = body["last_name"]
    db.commit()
    db.refresh(user)
    return {"detail": "Profile updated", "first_name": user.first_name, "last_name": user.last_name}


# ═══════════════════════════════════════════════
#  DOCUMENT PREVIEW (PPTX → HTML)
# ═══════════════════════════════════════════════

@app.get("/preview/{module_id}")
def preview_document(module_id: int, db: DB):
    """Convert uploaded documents (PPTX, etc.) to viewable HTML."""
    from fastapi.responses import HTMLResponse
    module = db.query(Module).filter(Module.id == module_id).first()
    if not module:
        raise HTTPException(status_code=404, detail="Module not found")

    file_path = module.file_path or module.url or ""
    if not file_path:
        return HTMLResponse("<p>No content available.</p>")

    # Resolve absolute path from the relative /uploads/... path
    abs_path = pathlib.Path(__file__).resolve().parent / file_path.lstrip("/")
    if not abs_path.exists():
        return HTMLResponse("<p>File not found on server.</p>")

    ext = abs_path.suffix.lower()

    if ext == ".pptx":
        try:
            from pptx import Presentation
            from pptx.util import Emu
            prs = Presentation(str(abs_path))
            slides_html = []
            for idx, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = para.text.strip()
                            if line:
                                # Detect headings (larger or bold text)
                                is_heading = False
                                for run in para.runs:
                                    if run.font.bold or (run.font.size and run.font.size > Emu(240000)):
                                        is_heading = True
                                        break
                                if is_heading:
                                    texts.append(f"<h3 style='margin:.5rem 0;color:#1a1a2e'>{_esc(line)}</h3>")
                                else:
                                    texts.append(f"<p style='margin:.3rem 0;line-height:1.6'>{_esc(line)}</p>")
                    elif shape.has_table:
                        table = shape.table
                        rows_html = ""
                        for row_idx, row in enumerate(table.rows):
                            cells = "".join(
                                f"<{'th' if row_idx == 0 else 'td'} style='border:1px solid #ddd;padding:.4rem .6rem'>{_esc(cell.text)}</{'th' if row_idx == 0 else 'td'}>"
                                for cell in row.cells
                            )
                            rows_html += f"<tr>{cells}</tr>"
                        texts.append(f"<table style='border-collapse:collapse;width:100%;margin:.5rem 0'>{rows_html}</table>")

                content = "".join(texts) if texts else "<p style='color:#999;font-style:italic'>Slide content (visual only)</p>"
                slides_html.append(
                    f"<div style='background:#fff;border-radius:10px;padding:1.5rem 2rem;margin-bottom:1.2rem;"
                    f"box-shadow:0 2px 8px rgba(0,0,0,.08);border:1px solid #e5e7eb'>"
                    f"<div style='font-size:.7rem;color:#6b7280;margin-bottom:.5rem;font-weight:600'>SLIDE {idx+1}</div>"
                    f"{content}</div>"
                )
            body = "".join(slides_html)
            return HTMLResponse(
                f"<div style='font-family:Inter,system-ui,sans-serif;max-width:800px;margin:0 auto;padding:1rem'>"
                f"<div style='font-size:.8rem;color:#6b7280;margin-bottom:1rem'>{len(prs.slides)} slides</div>"
                f"{body}</div>"
            )
        except Exception as e:
            return HTMLResponse(f"<p style='color:red'>Error reading PPTX: {_esc(str(e))}</p>")

    elif ext == ".docx":
        try:
            import docx
            doc = docx.Document(str(abs_path))
            paras = []
            for p in doc.paragraphs:
                text = p.text.strip()
                if not text:
                    continue
                if p.style.name.startswith("Heading"):
                    level = p.style.name[-1] if p.style.name[-1].isdigit() else "3"
                    paras.append(f"<h{level} style='margin:.8rem 0 .3rem'>{_esc(text)}</h{level}>")
                else:
                    paras.append(f"<p style='margin:.3rem 0;line-height:1.7'>{_esc(text)}</p>")
            return HTMLResponse(
                f"<div style='font-family:Inter,system-ui,sans-serif;max-width:780px;margin:0 auto;padding:1.5rem'>"
                f"{''.join(paras)}</div>"
            )
        except Exception as e:
            return HTMLResponse(f"<p style='color:red'>Error reading DOCX: {_esc(str(e))}</p>")

    else:
        # Plain text / md / html
        try:
            text = abs_path.read_text(encoding="utf-8", errors="replace")
            return HTMLResponse(
                f"<pre style='font-family:inherit;white-space:pre-wrap;line-height:1.7;padding:1.5rem;"
                f"max-width:780px;margin:0 auto'>{_esc(text)}</pre>"
            )
        except Exception:
            return HTMLResponse("<p>Could not read file.</p>")


def _esc(s: str) -> str:
    """HTML-escape a string."""
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")


# ═══════════════════════════════════════════════
#  Health Check
# ═══════════════════════════════════════════════

@app.get("/health")
def health_check():
    return {"status": "ok", "timestamp": datetime.now(timezone.utc).isoformat()}


# ═══════════════════════════════════════════════
#  Serve Frontend + Uploads Static Files
# ═══════════════════════════════════════════════

_FRONTEND_DIR = pathlib.Path(__file__).resolve().parent.parent / "frontend"
_CERT_DIR = pathlib.Path(__file__).resolve().parent / "certificates"
_CERT_DIR.mkdir(exist_ok=True)
app.mount("/uploads", StaticFiles(directory=str(UPLOAD_DIR)), name="uploads")
app.mount("/certificates", StaticFiles(directory=str(_CERT_DIR)), name="certificates")
if _FRONTEND_DIR.exists():
    app.mount("/", StaticFiles(directory=str(_FRONTEND_DIR), html=True), name="frontend")

# Seed admin on import
_seed_admin()
